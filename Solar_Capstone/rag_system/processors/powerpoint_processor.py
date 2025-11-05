"""PowerPoint Processor"""
from pathlib import Path
from typing import Optional
from .base_processor import BaseProcessor

try:
 from pptx import Presentation
 PPTX_AVAILABLE = True
except ImportError:
 PPTX_AVAILABLE = False

class PowerPointProcessor(BaseProcessor):
 def extract_text(self, file_path: Path) -> Optional[str]:
 if not PPTX_AVAILABLE:
 return "PowerPoint processing not available - install python-pptx"
 
 try:
 prs = Presentation(file_path)
 text = ""
 for slide in prs.slides:
 for shape in slide.shapes:
 if hasattr(shape, "text"):
 text += shape.text + "\n"
 return text
 except Exception as e:
 print(f" PowerPoint processing error: {e}")
 return None

